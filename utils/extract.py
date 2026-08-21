"""
utils/extract.py
Extracts raw text from uploaded notes (PDF, DOCX, PPTX, TXT).
Each function returns a list of (page_number, text) tuples so downstream
RAG code can cite a source page number.
"""

import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def extract_pdf(filepath):
    """Return list of (page_number, text) from a PDF file."""
    reader = PdfReader(filepath)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append((i, text))
    return pages


def extract_docx(filepath):
    """Return list of (page_number, text) from a DOCX file.
    Word has no real page breaks in the XML, so we treat each
    paragraph block of ~40 lines as one 'page' for citation purposes.
    """
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    pages = []
    chunk_size = 40
    for i in range(0, len(paragraphs), chunk_size):
        chunk = "\n".join(paragraphs[i:i + chunk_size])
        pages.append((i // chunk_size + 1, chunk))
    return pages if pages else [(1, "")]


def extract_pptx(filepath):
    """Return list of (slide_number, text) from a PPTX file."""
    prs = Presentation(filepath)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        pages.append((i, "\n".join(texts)))
    return pages


def extract_txt(filepath):
    """Return list of (page_number, text) from a plain text file.
    Splits into ~40-line chunks so citations behave consistently.
    """
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        lines = f.readlines()
    pages = []
    chunk_size = 40
    for i in range(0, len(lines), chunk_size):
        chunk = "".join(lines[i:i + chunk_size])
        pages.append((i // chunk_size + 1, chunk))
    return pages if pages else [(1, "")]


def extract_text(filepath):
    """Dispatch based on file extension. Returns list of (page_number, text)."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext in (".ppt", ".pptx"):
        return extract_pptx(filepath)
    elif ext == ".txt":
        return extract_txt(filepath)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def full_text(filepath):
    """Convenience helper: return the whole document as one string."""
    pages = extract_text(filepath)
    return "\n\n".join(text for _, text in pages)
